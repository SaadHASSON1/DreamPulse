from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Colors based on app theme
BG_COLOR = RGBColor(0x1A, 0x1B, 0x2E) # Deep Night Blue
TITLE_COLOR = RGBColor(0xE8, 0xE8, 0xF0) # Pearl White
TEXT_COLOR = RGBColor(0x90, 0xCA, 0xF9) # Soft Glow

prs = Presentation()

def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "DreamPulse"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
title.text_frame.paragraphs[0].font.bold = True

subtitle.text = "Akıllı Wear OS Uyku Alarmı\nUykunuzu optimize edin, dinç uyanın!"
subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

# Slide 2: Problem
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
title = slide.shapes.title
title.text = "Sorun Nedir?"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Yeterince uyumanıza rağmen yorgun mu uyanıyorsunuz?"
p = tf.add_paragraph()
p.text = "Geleneksel alarmlar sabit bir saate ayarlıdır."
p.level = 1
p = tf.add_paragraph()
p.text = "Ne zaman uykuya daldığınızı bilmezler."
p.level = 1
p = tf.add_paragraph()
p.text = "Derin uykunun ortasında uyanmak, gün boyu süren yorgunluğa neden olur."
p.level = 1

for paragraph in tf.paragraphs:
    paragraph.font.color.rgb = TEXT_COLOR

# Slide 3: Solution
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
title = slide.shapes.title
title.text = "Çözüm: DreamPulse"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "DreamPulse, saatinizi kullanarak uykuya daldığınız anı tespit eder."
p = tf.add_paragraph()
p.text = "Sabit saat yerine, hedeflediğiniz 'uyku süresine' odaklanır."
p.level = 1
p = tf.add_paragraph()
p.text = "Sadece uykuya daldıktan sonra geri sayıma başlar."
p.level = 1
p = tf.add_paragraph()
p.text = "Telefondan tamamen bağımsız çalışır."
p.level = 1

for paragraph in tf.paragraphs:
    paragraph.font.color.rgb = TEXT_COLOR

# Slide 4: How it works
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
title = slide.shapes.title
title.text = "Nasıl Çalışır?"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Kullanımı çok basittir:"
p = tf.add_paragraph()
p.text = "1. Uyumak istediğiniz süreyi seçin (Örn: 7.5 saat)."
p.level = 1
p = tf.add_paragraph()
p.text = "2. Saatinizi takın ve uyuyun."
p.level = 1
p = tf.add_paragraph()
p.text = "3. DreamPulse, kalp atış hızınızı analiz ederek uyku anınızı bulur."
p.level = 1
p = tf.add_paragraph()
p.text = "4. Hedef süreniz dolduğunda sizi uyandırır."
p.level = 1

for paragraph in tf.paragraphs:
    paragraph.font.color.rgb = TEXT_COLOR

# Slide 5: Call to action
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
title = slide.shapes.title
title.text = "Beta Testine Katılın!"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Google Play'de genel yayına geçmek için test kullanıcılarına ihtiyacımız var."
p = tf.add_paragraph()
p.text = "Google Play politikası gereği 12 test kullanıcısına ulaşmalıyız."
p.level = 1
p = tf.add_paragraph()
p.text = "Wear OS saat sahibiyseniz (Samsung Galaxy Watch, Pixel Watch vb.) sadece 1 gece deneyerek destek olabilirsiniz."
p.level = 1

for paragraph in tf.paragraphs:
    paragraph.font.color.rgb = TEXT_COLOR

# Slide 6: QR Code
slide_layout = prs.slide_layouts[5] # Title only
slide = prs.slides.add_slide(slide_layout)
set_dark_bg(slide)
title = slide.shapes.title
title.text = "Başlamak İçin QR Kodu Tarayın"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

img_path = r"C:\Users\saadh\.gemini\antigravity\brain\f3c18d9d-f4be-43b8-9aca-1422184a8133\dreampulse_qr_code_1778490183070.png"
if os.path.exists(img_path):
    left = Inches(3.0)
    top = Inches(2.0)
    height = Inches(4.5)
    slide.shapes.add_picture(img_path, left, top, height=height)

prs.save("DreamPulse_Sunum_TR.pptx")
print("Presentation created successfully!")
